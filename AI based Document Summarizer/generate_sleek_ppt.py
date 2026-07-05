from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_beautiful_ppt(output_path):
    prs = Presentation()
    
    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
        bg.line.fill.background()
        
        # Color bar
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.4), Inches(2.5), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(99, 102, 241)
        accent.line.fill.background()

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Segoe UI'

        # Subtitle
        subBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(2))
        stf = subBox.text_frame
        stf.word_wrap = True
        for i, line in enumerate(subtitle_text.split('\n')):
            p2 = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
            p2.text = line
            p2.font.size = Pt(22)
            p2.font.color.rgb = RGBColor(148, 163, 184)
            p2.font.name = 'Segoe UI'

    def add_content_slide(title_text, bullets, placeholder=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
        bg.line.fill.background()
        
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(99, 102, 241)
        top_bar.line.fill.background()

        tBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = tBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Segoe UI'

        box_width = 8.5 if not placeholder else 5.2
        cBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(box_width), Inches(5))
        ctf = cBox.text_frame
        ctf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
            p.text = bullet.strip()
            p.font.size = Pt(21)
            p.font.name = 'Segoe UI'
            p.font.color.rgb = RGBColor(241, 245, 249)
            
            if bullet.startswith("    "):
                p.level = 1
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(148, 163, 184)
        
        if placeholder:
            ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(2.0), Inches(3.6), Inches(4.2))
            ph.fill.solid()
            ph.fill.fore_color.rgb = RGBColor(30, 41, 59)
            ph.line.color.rgb = RGBColor(99, 102, 241)
            ph.line.width = Pt(2)
            ptf = ph.text_frame
            ptf.word_wrap = True
            for _ in range(4): ptf.add_paragraph()
            pp = ptf.paragraphs[4]
            pp.text = f"[ DELETE ME & \nINSERT {placeholder} GFX ]"
            pp.font.bold = True
            pp.font.size = Pt(16)
            pp.font.color.rgb = RGBColor(148, 163, 184)
            from pptx.enum.text import PP_ALIGN
            pp.alignment = PP_ALIGN.CENTER

    # Building Slides
    add_title_slide(
        "Heart Disease Risk Assessment\nusing Predictive Analytics", 
        "A Comparative Study of Machine Learning Algorithms\nCourse: BCA (Final Year) - Continuous Assessment\nUniversity: Lovely Professional University (LPU)"
    )

    add_content_slide("Introduction & Objective", [
        "The Problem: Cardiovascular diseases (CVDs) are the leading cause of death globally. Early detection is a critical challenge.",
        "The Goal: To develop a Machine Learning pipeline that predicts heart disease risk based on clinical parameters.",
        "Tech Stack:",
        "    Python (Google Colab Environment)",
        "    Pandas & NumPy (Data Handling)",
        "    Matplotlib & Seaborn (Visualization)",
        "    Scikit-Learn (Machine Learning Models)"
    ])

    add_content_slide("Dataset Description", [
        "Source: UCI Machine Learning Repository (Cleveland Dataset).",
        "Scope: 303 patient samples with 14 clinical attributes.",
        "Key Attributes:",
        "    trestbps: Resting blood pressure (mm Hg).",
        "    chol: Serum cholesterol (mg/dl).",
        "    thalach: Maximum heart rate achieved.",
        "    oldpeak: ST depression induced by exercise.",
        "    ca: Number of major vessels colored by flourosopy."
    ])

    add_content_slide("Data Preprocessing", [
        "Feature Selection: Removed irrelevant metadata such as ID to prevent noise.",
        "Target Simplification: Transformed standard attributes into a binary target:",
        "    0: No Heart Disease",
        "    1: Presence of Heart Disease",
        "Categorical Encoding: Applied One-Hot Encoding to convert text-based features into numeric formats.",
        "Handling Missing Data: Employed Median Imputation to prevent statistical bias."
    ])

    add_content_slide("Feature Scaling", [
        "Method: StandardScaler (Z-score Normalization).",
        "Rationale: Health data has vastly different scales (Age: 20-80 vs. Cholesterol).",
        "Importance: Essential for distance-based algorithms like SVM.",
        "    Guarantees no single high-value feature dominates the model's decision-making process."
    ])

    add_content_slide("Comparative Study of Algorithms", [
        "Evaluated approaches:",
        "    1. Logistic Regression: Baseline model.",
        "    2. SVM (Support Vector Machine): RBF kernel.",
        "    3. Random Forest Classifier: Ensemble method.",
        "Results:",
        "    Logistic Regression: ~80.4%",
        "    SVM Accuracy: ~83.7%",
        "    Random Forest Accuracy: 86.41%"
    ], "COMPARISON")

    add_content_slide("Performance (Random Forest)", [
        "Accuracy: 86.41% on the test dataset.",
        "Confusion Matrix Analysis:",
        "    True Positives/Negatives: High count.",
        "    False Negatives: Minimized to ensure safety.",
        "Recall (Sensitivity): Prioritized because missing a",
        "    high-risk patient is a critical error."
    ], "CONFUSION MATRIX")

    add_content_slide("Conclusion", [
        "The project successfully demonstrates the utility of Predictive Analytics in medical diagnostics.",
        "Random Forest proved to be robust algorithm for the dataset.",
        "Outcome: A trained model capable of providing a decision-support metric for clinicians."
    ])

    add_content_slide("Future Scope & Q&A", [
        "Deployment: Integrating the model into a web dashboard for real-time assessments.",
        "Data Expansion: Training with larger, more diverse datasets to improve sensitivity.",
        "Questions?"
    ])

    # Ensure output directory exists and save
    prs.save(output_path)

if __name__ == '__main__':
    create_beautiful_ppt(r"C:\Users\aarya\Desktop\Heart_Disease_Premium_Deck.pptx")
    print("Beautiful Premium PPT created!")
