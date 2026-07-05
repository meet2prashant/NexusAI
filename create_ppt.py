from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Heart Disease Risk Assessment using Predictive Analytics"
subtitle.text = "A Comparative Study of Machine Learning Algorithms\nCourse: BCA (Final Year) - Continuous Assessment\nUniversity: Lovely Professional University (LPU)"

# Helper function
def add_slide(title_text, content_bullets):
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    sl.shapes.title.text = title_text
    tf = sl.placeholders[1].text_frame
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            if bullet.startswith("    "):
                p.level = 1
                p.text = bullet.strip()
    return sl

add_slide("Introduction & Objective", [
    "The Problem: Cardiovascular diseases (CVDs) are the leading cause of death globally. Early detection is a critical challenge in clinical settings.",
    "The Goal: To develop a Machine Learning pipeline that predicts heart disease risk based on multi-dimensional clinical parameters.",
    "Tech Stack:",
    "    Python (Google Colab Environment)",
    "    Pandas & NumPy (Data Handling)",
    "    Matplotlib & Seaborn (Visualization)",
    "    Scikit-Learn (Machine Learning Models)"
])

add_slide("Dataset Description", [
    "Source: UCI Machine Learning Repository (Cleveland Dataset).",
    "Scope: 303 patient samples with 14 clinical attributes.",
    "Key Attributes:",
    "    trestbps: Resting blood pressure (mm Hg).",
    "    chol: Serum cholesterol (mg/dl).",
    "    thalach: Maximum heart rate achieved.",
    "    oldpeak: ST depression induced by exercise relative to rest.",
    "    ca: Number of major vessels colored by flourosopy."
])

add_slide("Data Preprocessing", [
    "Feature Selection: Removed irrelevant metadata such as id and dataset to prevent noise.",
    "Target Simplification: Transformed the num attribute (levels 0-4) into a binary target:",
    "    0: No Heart Disease",
    "    1: Presence of Heart Disease",
    "Categorical Encoding: Applied One-Hot Encoding (pd.get_dummies) to convert text-based features into numeric formats.",
    "Handling Missing Data: Employed Median Imputation to fill missing entries without introducing statistical bias."
])

add_slide("Feature Scaling", [
    "Method: StandardScaler (Z-score Normalization).",
    "Rationale: Health data has vastly different scales (Age: 20-80 vs. Cholesterol: 150-400).",
    "Importance: This step is essential for distance-based algorithms like SVM to ensure no single high-value feature dominates the model's decision-making process."
])

add_slide("Comparative Study of Algorithms", [
    "Evaluated approaches:",
    "    1. Logistic Regression: Used as a baseline linear model.",
    "    2. Support Vector Machine (SVM): Used with an RBF kernel to capture non-linear relationships.",
    "    3. Random Forest Classifier: An ensemble method utilizing multiple decision trees.",
    "Results:",
    "    Logistic Regression Accuracy: ~80.4%",
    "    SVM Accuracy: ~83.7%",
    "    Random Forest Accuracy: 86.41% (Final Selected Model)",
    "",
    "    [INSERT MODEL COMPARISON SCREENSHOT HERE]"
])

add_slide("Performance Evaluation (Random Forest)", [
    "Accuracy: 86.41% on the test dataset.",
    "Confusion Matrix Analysis:",
    "    True Positives/Negatives: High count, showing strong generalization.",
    "    False Negatives: Minimized to ensure sick patients are not incorrectly cleared.",
    "Recall (Sensitivity): Prioritized because missing a high-risk patient is a critical error in healthcare informatics.",
    "",
    "    [INSERT CONFUSION MATRIX SCREENSHOT HERE]"
])

add_slide("Conclusion", [
    "The project successfully demonstrates the utility of Predictive Analytics in medical diagnostics.",
    "Random Forest proved to be the most robust algorithm for this specific clinical dataset.",
    "Outcome: A trained model capable of providing a decision-support metric for clinicians based on standard health parameters."
])

add_slide("Future Scope & Q&A", [
    "Deployment: Integrating the model into a web dashboard (using Flask or Streamlit) for real-time risk assessment.",
    "Data Expansion: Training with larger, more diverse datasets to improve model sensitivity across different demographics.",
    "Questions?"
])

prs.save(r"C:\Users\aarya\Desktop\Heart_Disease_Risk_Assessment.pptx")
print("Presentation generated successfully!")
